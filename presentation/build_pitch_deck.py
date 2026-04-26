"""
Generates AgenticFrmk VC Pitch Deck — AgenticFrmk-VC-Pitch.pptx
Run: python3.14 build_pitch_deck.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette ────────────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x08, 0x0b, 0x10)   # page background
BG_CARD    = RGBColor(0x0d, 0x11, 0x17)   # card/box fill
BG_MID     = RGBColor(0x16, 0x1b, 0x22)   # subtle panel
BORDER     = RGBColor(0x1c, 0x23, 0x33)   # card border
PURPLE     = RGBColor(0x7c, 0x3a, 0xed)   # primary accent
PURPLE_LT  = RGBColor(0xa7, 0x8b, 0xfa)   # light purple text
AMBER      = RGBColor(0xf5, 0x9e, 0x0b)   # warning accent
GREEN      = RGBColor(0x22, 0xc5, 0x5e)   # success accent
TEXT_WHITE = RGBColor(0xe6, 0xed, 0xf3)   # body text
TEXT_DIM   = RGBColor(0x8b, 0x94, 0x9e)   # secondary text
TEXT_DIMMER= RGBColor(0x48, 0x4f, 0x58)   # muted

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ── Low-level helpers ──────────────────────────────────────────────────────────

def add_slide(prs: Presentation) -> object:
    blank = prs.slide_layouts[6]   # completely blank layout
    return prs.slides.add_slide(blank)


def fill_slide(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h,
        fill=None, border_color=None, border_pt=0.75, radius=0) -> object:
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.line.width = Pt(border_pt) if border_color else 0
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()

    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def txt(slide, text: str, x, y, w, h,
        size=18, bold=False, color=TEXT_WHITE,
        align=PP_ALIGN.LEFT, wrap=True) -> object:
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Inter"
    return txb


def multiline(slide, lines: list, x, y, w, h,
              size=14, bold=False, color=TEXT_WHITE,
              line_color=None, align=PP_ALIGN.LEFT,
              spacing_pt=2) -> object:
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (line_text, line_bold, line_size, line_col) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing_pt)
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(line_size or size)
        run.font.bold = line_bold if line_bold is not None else bold
        run.font.color.rgb = line_col or color
        run.font.name = "Inter"
    return txb


def pill(slide, text: str, x, y, color: RGBColor, bg: RGBColor,
         size=9, w=1.3, h=0.28) -> None:
    b = box(slide, x, y, w, h, fill=bg, border_color=None)
    t = txt(slide, text, x + 0.05, y + 0.02, w - 0.1, h - 0.04,
            size=size, bold=True, color=color, align=PP_ALIGN.CENTER)


def divider(slide, y, color=BORDER) -> None:
    line = slide.shapes.add_connector(1, Inches(0.4), Inches(y), Inches(12.93), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(0.5)


def accent_bar(slide, x, y, h=0.5, w=0.06, color=PURPLE) -> None:
    b = box(slide, x, y, w, h, fill=color)


# ── SLIDE BUILDERS ─────────────────────────────────────────────────────────────

def slide_title(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    # Gradient-ish decorative bars
    box(sl, 0, 0, 0.5, 7.5, fill=RGBColor(0x12, 0x06, 0x2a))
    box(sl, 12.83, 0, 0.5, 7.5, fill=RGBColor(0x12, 0x06, 0x2a))

    # Purple accent glow circle (approximated as a small box)
    box(sl, 5.8, 0.4, 1.8, 1.0, fill=RGBColor(0x1a, 0x0a, 0x38))

    txt(sl, "AgenticFrmk", 1.5, 1.5, 10, 1.2,
        size=54, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "The production-grade agentic infrastructure layer\nthat enterprises actually deploy.",
        1.5, 2.85, 10, 1.1, size=22, color=PURPLE_LT, align=PP_ALIGN.CENTER)

    divider(sl, 4.15)

    # 5 service pills
    services = [
        ("AuthService", RGBColor(0x22, 0xc5, 0x5e), RGBColor(0x06, 0x2a, 0x14)),
        ("AgentGateway", PURPLE_LT, RGBColor(0x1e, 0x10, 0x45)),
        ("AgentCore", RGBColor(0x38, 0xbd, 0xf8), RGBColor(0x07, 0x21, 0x2e)),
        ("RegistryService", AMBER, RGBColor(0x2c, 0x1a, 0x03)),
        ("SREDemo", RGBColor(0xfb, 0x71, 0x85), RGBColor(0x2c, 0x08, 0x14)),
    ]
    start_x = 1.3
    for name, col, bg in services:
        pill(sl, name, start_x, 4.35, col, bg, size=10, w=2.1, h=0.35)
        start_x += 2.15

    txt(sl, "Agentic Seattle Hackathon · 2026", 1.5, 5.0, 10, 0.5,
        size=12, color=TEXT_DIMMER, align=PP_ALIGN.CENTER)

    txt(sl, "mauttaram@gmail.com", 1.5, 5.5, 10, 0.5,
        size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    txt(sl, "github.com/AgenticFrmk", 1.5, 5.9, 10, 0.5,
        size=12, color=PURPLE_LT, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4, h=0.6)
    txt(sl, "The Problem", 0.72, 0.38, 8, 0.65, size=30, bold=True, color=TEXT_WHITE)
    txt(sl, "02", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    # Big stat box
    box(sl, 0.5, 1.2, 4.0, 2.2, fill=RGBColor(0x12, 0x06, 0x2a), border_color=PURPLE)
    txt(sl, "$5,600", 0.6, 1.35, 3.8, 1.0, size=52, bold=True, color=PURPLE_LT, align=PP_ALIGN.CENTER)
    txt(sl, "per minute of P1 downtime", 0.6, 2.35, 3.8, 0.55,
        size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    box(sl, 4.7, 1.2, 3.8, 2.2, fill=BG_CARD, border_color=BORDER)
    txt(sl, "45–90 min", 4.8, 1.35, 3.6, 0.9, size=40, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    txt(sl, "manual triage per P1 incident", 4.8, 2.25, 3.6, 0.65,
        size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    box(sl, 8.7, 1.2, 4.1, 2.2, fill=BG_CARD, border_color=BORDER)
    txt(sl, "$250K–$500K", 8.8, 1.35, 3.9, 0.9, size=28, bold=True, color=RGBColor(0xf8, 0x71, 0x71), align=PP_ALIGN.CENTER)
    txt(sl, "cost per incident (labour + downtime)", 8.8, 2.25, 3.9, 0.65,
        size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    divider(sl, 3.6)

    txt(sl, "The obvious answer: an AI agent.   The problem: every agent framework built today is a demo.",
        0.5, 3.7, 12.3, 0.6, size=14, color=TEXT_DIM)

    # Comparison table header
    col_positions = [0.5, 5.3]
    headers = ["What demo agents do", "What production needs"]
    header_colors = [RGBColor(0xf8, 0x71, 0x71), GREEN]
    for i, (hdr, col) in enumerate(zip(headers, header_colors)):
        box(sl, col_positions[i], 4.45, 4.6, 0.35, fill=BG_MID)
        txt(sl, hdr, col_positions[i] + 0.15, 4.48, 4.3, 0.32,
            size=10, bold=True, color=col)

    rows = [
        ("Single LLM call", "9-node LangGraph StateGraph"),
        ("Sequential steps", "Auto DAG-based parallel execution"),
        ("No crash recovery", "Claim/lease + auto re-claim on crash"),
        ("Black-box reasoning", "Explicit CoT trace surfaced to user"),
        ("No auth", "RS256 JWT + session identity binding"),
        ("No memory management", "Context budget + sliding-window compaction"),
    ]
    row_y = 4.88
    for left, right in rows:
        txt(sl, f"✕  {left}", 0.5, row_y, 4.6, 0.3, size=10,
            color=RGBColor(0xf8, 0x71, 0x71))
        txt(sl, f"✓  {right}", 5.3, row_y, 7.5, 0.3, size=10, color=GREEN)
        row_y += 0.3


def slide_solution_overview(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4)
    txt(sl, "The Solution: AgenticFrmk Platform", 0.72, 0.38, 10, 0.6,
        size=28, bold=True, color=TEXT_WHITE)
    txt(sl, "03", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    # 5 service cards in a row
    services = [
        ("AuthService", "RS256 JWT\nissuance\nJWKS rotation",
         GREEN, RGBColor(0x06, 0x2a, 0x14)),
        ("AgentGateway", "JWT verify\nSession bind\nContext budget\nCompaction",
         PURPLE_LT, RGBColor(0x1e, 0x10, 0x45)),
        ("AgentCore", "LangGraph\n9-node graph\nDAG executor\nHITL gates\n121 tests",
         RGBColor(0x38, 0xbd, 0xf8), RGBColor(0x07, 0x21, 0x2e)),
        ("RegistryService", "Schema Registry\nTool Registry\nPlaybook Rules\nVersioned API",
         AMBER, RGBColor(0x2c, 0x1a, 0x03)),
        ("SREDemo", "React 18 UI\nFastAPI SSE\nPlan History\n4 scenarios",
         RGBColor(0xfb, 0x71, 0x85), RGBColor(0x2c, 0x08, 0x14)),
    ]

    card_x = 0.45
    for name, body, col, bg in services:
        box(sl, card_x, 1.15, 2.4, 3.8, fill=bg, border_color=col)
        txt(sl, name, card_x + 0.15, 1.25, 2.1, 0.45, size=13, bold=True, color=col)
        divider_y = 1.78
        b = box(sl, card_x + 0.05, 1.72, 2.3, 0.04, fill=col)
        txt(sl, body, card_x + 0.15, 1.82, 2.1, 3.0, size=10.5, color=TEXT_DIM)
        card_x += 2.52

    # Bottom: key claim
    box(sl, 0.5, 5.3, 12.33, 1.8, fill=BG_CARD, border_color=BORDER)
    txt(sl, "\"These aren't edge cases. They're the default failure mode of every agent deployed today.\"",
        0.8, 5.5, 11.8, 0.75, size=15, bold=True, color=PURPLE_LT, align=PP_ALIGN.CENTER)
    txt(sl, "AgenticFrmk is the infrastructure layer that closes every one of these gaps — in production, not in a notebook.",
        0.8, 6.1, 11.8, 0.75, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)


def slide_gateway(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4)
    txt(sl, "AgentGateway", 0.72, 0.38, 6, 0.6, size=28, bold=True, color=TEXT_WHITE)
    txt(sl, "The production entry point nobody else built", 0.72, 0.95, 8, 0.45,
        size=15, color=PURPLE_LT)
    txt(sl, "04", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    features = [
        ("RS256 JWT Verification",
         "JWKS-rotated keypairs from AuthService. Only authenticated, identity-bound users can invoke agents. Zero hardcoded secrets.",
         PURPLE_LT),
        ("Session Identity Binding",
         "Every thread is locked to the originating user's JWT sub. No session hijacking. Resuming a thread as a different user returns 403.",
         RGBColor(0x38, 0xbd, 0xf8)),
        ("Real-Time Context Budget",
         "Token counting per session. Configurable warning threshold (70%) and compaction threshold (80%). Surfaced live to the UI.",
         AMBER),
        ("Sliding-Window Compaction",
         "When context fills: evict 30% oldest messages → LLM summarises → entity fidelity check → proceed. User sees it happen.",
         GREEN),
        ("Model Registry",
         "GET /models serves available models from config. The UI never hardcodes model names. New models ship without frontend changes.",
         RGBColor(0xfb, 0x71, 0x85)),
    ]

    card_y = 1.55
    for title, body, col in features:
        box(sl, 0.5, card_y, 12.33, 0.9, fill=BG_CARD, border_color=BORDER)
        box(sl, 0.5, card_y, 0.08, 0.9, fill=col)
        txt(sl, title, 0.72, card_y + 0.06, 3.5, 0.35, size=11, bold=True, color=col)
        txt(sl, body, 4.3, card_y + 0.06, 8.3, 0.75, size=10.5, color=TEXT_DIM)
        card_y += 1.0

    txt(sl, "This is the difference between a demo and infrastructure a Fortune 500 CISO will sign off on.",
        0.5, 6.85, 12.33, 0.45, size=12, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)


def slide_plan_memory(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4, color=PURPLE)
    txt(sl, "Plan Memory + Few-Shot Learning", 0.72, 0.38, 9, 0.6,
        size=28, bold=True, color=TEXT_WHITE)
    txt(sl, "The agent that gets smarter with every incident", 0.72, 0.95, 9, 0.45,
        size=15, color=PURPLE_LT)
    txt(sl, "05", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    # Left: how it works
    box(sl, 0.5, 1.55, 6.0, 4.3, fill=BG_CARD, border_color=BORDER)
    txt(sl, "How It Works", 0.72, 1.65, 5.6, 0.45, size=13, bold=True, color=PURPLE_LT)

    steps = [
        ("1", "SRE resolves a K8s crashloop — 9 steps, 287 seconds"),
        ("2", "Resolution stored: action, domain, steps, duration, outcome"),
        ("3", "Next K8s incident → agent searches Plan History"),
        ("4", "Match found → inject as few-shot example into plan prompt"),
        ("5", "Agent produces plan faster, aligned to proven patterns"),
        ("6", '"Few-shot used by agent" badge appears in UI history panel'),
    ]
    step_y = 2.15
    for num, desc in steps:
        box(sl, 0.65, step_y, 0.38, 0.38, fill=PURPLE, border_color=None)
        txt(sl, num, 0.65, step_y + 0.02, 0.38, 0.35,
            size=11, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, desc, 1.15, step_y + 0.04, 5.15, 0.32, size=10.5, color=TEXT_DIM)
        step_y += 0.52

    # Right: moat
    box(sl, 6.8, 1.55, 6.03, 4.3, fill=RGBColor(0x1e, 0x10, 0x45), border_color=PURPLE)
    txt(sl, "The Compounding Moat", 7.0, 1.65, 5.6, 0.45,
        size=13, bold=True, color=PURPLE_LT)

    moat_points = [
        "Year 1: agent learns your K8s patterns",
        "Year 2: 500+ incidents encoded → 40% faster resolution",
        "Year 3: new engineer joins → inherits all institutional knowledge",
        "",
        "Competitors starting fresh have NONE of this history.",
        "Switching cost grows with every incident resolved.",
    ]
    moat_y = 2.2
    for pt in moat_points:
        col = PURPLE_LT if pt.startswith("Competitors") or pt.startswith("Switching") else TEXT_DIM
        bold = pt.startswith("Competitors") or pt.startswith("Switching")
        txt(sl, pt, 7.0, moat_y, 5.6, 0.45, size=11, bold=bold, color=col)
        moat_y += 0.5

    # Bottom stat
    box(sl, 0.5, 6.05, 12.33, 0.98, fill=BG_MID, border_color=BORDER)
    txt(sl, "No retraining. No manual curation. The agent improves automatically — just by running.",
        0.7, 6.2, 11.9, 0.65, size=14, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)


def slide_hitl_registry(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4, color=AMBER)
    txt(sl, "HITL + RegistryService", 0.72, 0.38, 8, 0.6,
        size=28, bold=True, color=TEXT_WHITE)
    txt(sl, "Hard gates + grounded reasoning = enterprise trust", 0.72, 0.95, 9, 0.45,
        size=14, color=AMBER)
    txt(sl, "06", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    # HITL section
    box(sl, 0.5, 1.5, 5.8, 4.6, fill=BG_CARD, border_color=BORDER)
    txt(sl, "Human-in-the-Loop (Hard Gates)", 0.7, 1.6, 5.4, 0.45,
        size=12, bold=True, color=AMBER)

    hitl_rows = [
        ("Plan Review Gate", "SRE sees full DAG before any execution"),
        ("CoT Confirm Gate", "Claude Opus walks the reasoning aloud"),
        ("Modify → Re-plan", "Feedback routes back to planner, not skipped"),
        ("Audit Trail", "Every approval: user identity + timestamp logged"),
        ("Compliance Ready", "PCI-DSS / SOC2 / EU AI Act evidence artifact"),
    ]
    row_y = 2.12
    for title, body in hitl_rows:
        txt(sl, f"⬡  {title}", 0.7, row_y, 2.8, 0.35, size=10.5, bold=True, color=AMBER)
        txt(sl, body, 3.5, row_y, 2.6, 0.35, size=10.5, color=TEXT_DIM)
        row_y += 0.58

    txt(sl, "The executor CANNOT start\nwithout explicit human approval.",
        0.7, 4.5, 5.4, 0.75, size=11, bold=True, color=TEXT_WHITE)

    # RegistryService section
    box(sl, 6.6, 1.5, 6.2, 4.6, fill=BG_CARD, border_color=BORDER)
    txt(sl, "RegistryService (Anti-Hallucination)", 6.8, 1.6, 5.8, 0.45,
        size=12, bold=True, color=RGBColor(0xf5, 0x9e, 0x0b))

    reg_rows = [
        ("Schema Registry", "Entity fields → extract_intent, extract_entities", RGBColor(0x38, 0xbd, 0xf8)),
        ("Tool Registry", "Tool names + contracts → plan node uses real names only", GREEN),
        ("Playbook Registry", "Domain ordering rules, hard constraints → validated post-plan", AMBER),
        ("Team-Owned", "POST /schemas, /tools, /playbooks — no code changes needed", PURPLE_LT),
        ("Versioned", "BACKWARD/FORWARD/FULL/BREAKING compat checking", RGBColor(0xfb, 0x71, 0x85)),
    ]
    row_y = 2.12
    for title, body, col in reg_rows:
        txt(sl, f"▸  {title}", 6.8, row_y, 2.8, 0.38, size=10.5, bold=True, color=col)
        txt(sl, body, 9.6, row_y, 3.0, 0.38, size=10, color=TEXT_DIM)
        row_y += 0.58

    txt(sl, "The LLM never guesses tool names.\nIt reasons from registered facts.",
        6.8, 4.5, 5.8, 0.75, size=11, bold=True, color=TEXT_WHITE)

    # Bottom
    box(sl, 0.5, 6.3, 12.33, 0.8, fill=BG_MID, border_color=BORDER)
    txt(sl, "Together: every consequential action is human-approved and grounded in versioned domain knowledge.",
        0.7, 6.45, 11.9, 0.5, size=12, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)


def slide_sre_demo_value(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4, color=GREEN)
    txt(sl, "SREDemo: Business Value in 5 Minutes", 0.72, 0.38, 10, 0.6,
        size=26, bold=True, color=TEXT_WHITE)
    txt(sl, "07", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    # Incident trigger
    box(sl, 0.5, 1.2, 12.33, 0.7, fill=BG_CARD, border_color=BORDER)
    txt(sl, '🔴  SRE types: "K8s pods are OOMKilled every 2 minutes in production — payment-service is down"',
        0.7, 1.32, 11.9, 0.45, size=12, color=TEXT_DIM)

    # Phase timeline
    phases = [
        ("Intent",    "Classifies domain, action, severity",          "~3s",  PURPLE_LT),
        ("Clarify",   "Asks one targeted question",                    "~5s",  AMBER),
        ("Entities",  "Extracts severity, service, error type",        "~3s",  PURPLE_LT),
        ("Plan",      "9-step DAG — few-shot from history",            "~8s",  GREEN),
        ("HITL Gate", "SRE reviews + approves full plan",              "Human",RGBColor(0xf8, 0x71, 0x71)),
        ("Execute",   "9 tools run in parallel dependency order",      "~45s", RGBColor(0x38, 0xbd, 0xf8)),
        ("Report",    "Root cause + remediation + follow-ups",         "~5s",  GREEN),
    ]

    col_w = 1.76
    col_x = 0.5
    for name, desc, time_val, col in phases:
        box(sl, col_x, 2.1, col_w - 0.08, 2.4, fill=BG_CARD, border_color=col)
        box(sl, col_x, 2.1, col_w - 0.08, 0.06, fill=col)
        txt(sl, name, col_x + 0.1, 2.2, col_w - 0.25, 0.38, size=11, bold=True, color=col)
        txt(sl, desc, col_x + 0.1, 2.62, col_w - 0.25, 0.85, size=9.5, color=TEXT_DIM)
        txt(sl, time_val, col_x + 0.1, 3.75, col_w - 0.25, 0.45,
            size=13, bold=True, color=col if name == "HITL Gate" else TEXT_DIM, align=PP_ALIGN.LEFT)
        col_x += col_w

    # Savings boxes
    box(sl, 0.5, 4.75, 5.9, 1.5, fill=RGBColor(0x06, 0x2a, 0x14), border_color=GREEN)
    txt(sl, "< 5 minutes", 0.7, 4.9, 5.5, 0.65, size=34, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(sl, "agent-assisted resolution (vs 45–90 min manual)", 0.7, 5.5, 5.5, 0.55,
        size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    box(sl, 6.6, 4.75, 6.23, 1.5, fill=RGBColor(0x12, 0x06, 0x2a), border_color=PURPLE)
    txt(sl, "$224K–$476K", 6.8, 4.9, 5.9, 0.65, size=30, bold=True, color=PURPLE_LT, align=PP_ALIGN.CENTER)
    txt(sl, "saved per P1 incident at $5,600/min", 6.8, 5.5, 5.9, 0.55,
        size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # 4 scenarios
    scenarios = [
        ("☸ K8s Crashloop", "OOMKilled · Clarification"),
        ("🐘 DB Pool", "Connection exhaustion"),
        ("🔌 VPN Flap", "IKE phase 2 mismatch"),
        ("🔒 SSL Expiry", "Auto-renewal failure"),
    ]
    scen_x = 0.5
    for name, sub in scenarios:
        box(sl, scen_x, 6.45, 3.0, 0.75, fill=BG_MID, border_color=BORDER)
        txt(sl, name, scen_x + 0.15, 6.52, 2.7, 0.3, size=10.5, bold=True, color=TEXT_WHITE)
        txt(sl, sub, scen_x + 0.15, 6.82, 2.7, 0.28, size=9, color=TEXT_DIM)
        scen_x += 3.1


def slide_architecture(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4, color=RGBColor(0x38, 0xbd, 0xf8))
    txt(sl, "Architecture", 0.72, 0.38, 8, 0.6, size=28, bold=True, color=TEXT_WHITE)
    txt(sl, "08", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    # Stack diagram as boxes with arrows
    layers = [
        ("Browser · React 18 + TypeScript",
         "EventSource SSE · fetch · Budget Gauge · Plan History · Compaction Banners",
         RGBColor(0x38, 0xbd, 0xf8), RGBColor(0x07, 0x21, 0x2e)),
        ("SREDemo Backend · FastAPI + Python 3.13",
         "asyncio.Queue per session · Script dispatch · Plan History store · SSE stream",
         RGBColor(0xfb, 0x71, 0x85), RGBColor(0x2c, 0x08, 0x14)),
        ("AgentGateway · FastAPI + Python 3.13",
         "RS256 JWT verify · Session bind · Context budget · Sliding-window compaction · Model registry",
         PURPLE_LT, RGBColor(0x1e, 0x10, 0x45)),
        ("AgentCore · LangGraph StateGraph",
         "9 nodes: intent → clarify → entities → plan → HITL → CoT → execute_step × N (parallel) → report\n"
         "Send API fan-out · interrupt() gates · MemorySaver / PostgresSaver · 121 tests",
         RGBColor(0x34, 0xd3, 0x99), RGBColor(0x06, 0x2a, 0x1e)),
    ]

    layer_y = 1.15
    for name, body, col, bg in layers:
        box(sl, 0.5, layer_y, 9.5, 1.05, fill=bg, border_color=col)
        txt(sl, name, 0.7, layer_y + 0.06, 9.1, 0.38, size=11, bold=True, color=col)
        txt(sl, body, 0.7, layer_y + 0.46, 9.1, 0.55, size=9.5, color=TEXT_DIM)
        # Arrow
        if layer_y < 4.3:
            txt(sl, "↓", 4.7, layer_y + 1.08, 0.5, 0.35, size=14, color=col, align=PP_ALIGN.CENTER)
        layer_y += 1.35

    # Sidebar: supporting services
    box(sl, 10.25, 1.15, 2.6, 2.3, fill=BG_CARD, border_color=BORDER)
    txt(sl, "AuthService", 10.4, 1.22, 2.3, 0.35, size=10, bold=True, color=GREEN)
    txt(sl, "RS256 key generation\nJWT issuance\nJWKS endpoint", 10.4, 1.58, 2.3, 0.75, size=9, color=TEXT_DIM)

    box(sl, 10.25, 3.6, 2.6, 2.3, fill=BG_CARD, border_color=BORDER)
    txt(sl, "RegistryService", 10.4, 3.67, 2.3, 0.35, size=10, bold=True, color=AMBER)
    txt(sl, "Schema Registry\nTool Registry\nPlaybook Registry\nVersioned + Compat-checked\n60s TTL client cache",
        10.4, 4.03, 2.3, 1.5, size=9, color=TEXT_DIM)

    txt(sl, "Claude Sonnet 4.6 (orchestration) · Claude Opus 4.7 (CoT validation)",
        0.5, 6.05, 9.5, 0.4, size=10, color=TEXT_DIM)

    # AgentCore node graph
    box(sl, 0.5, 6.55, 12.33, 0.75, fill=BG_MID, border_color=BORDER)
    txt(sl, "extract_intent  →  [clarify]  →  extract_entities  →  plan  →  HITL  →  validate_cot  →  executor_router  →→  execute_step × N  →  report",
        0.7, 6.65, 11.9, 0.5, size=9, color=PURPLE_LT, align=PP_ALIGN.CENTER)


def slide_moat(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4, color=PURPLE)
    txt(sl, "The Moat", 0.72, 0.38, 8, 0.6, size=28, bold=True, color=TEXT_WHITE)
    txt(sl, "Compounding advantages that widen over time", 0.72, 0.95, 9, 0.45,
        size=14, color=PURPLE_LT)
    txt(sl, "09", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    moats = [
        ("📚 Plan Memory Compounds",
         "Every incident resolution → few-shot library. The agent improves on your specific infrastructure without retraining. Competitors can't copy years of encoded patterns.",
         PURPLE_LT),
        ("🗂 RegistryService Compounds",
         "Every domain team registers their schemas, tools, playbooks. New LLM sees all of it immediately. Accumulated domain knowledge is not copyable.",
         AMBER),
        ("✅ HITL as Compliance Artifact",
         "Every change: human-approved, user-identity-timestamped. PCI-DSS / SOC2 / EU AI Act evidence is automatic, not retrofitted.",
         GREEN),
        ("🔐 Gateway as Enterprise Gate",
         "RS256 JWT + session binding + context budget = the layer a Fortune 500 CISO actually signs off on. No other framework ships this.",
         RGBColor(0x38, 0xbd, 0xf8)),
        ("🕸 DAG Safety",
         "Partial broken configs are structurally impossible. Dependencies enforced before execution. Playbook hard rules reject bad plans before any tool runs.",
         RGBColor(0xfb, 0x71, 0x85)),
        ("📏 Horizontal Scale",
         "DispatchAdapter pattern: swap Postgres → SQS → Pub/Sub with one env var, zero code changes. Scales to Fortune 500 without re-architecture.",
         RGBColor(0x34, 0xd3, 0x99)),
    ]

    col = 0
    row = 0
    for title, body, color in moats:
        x = 0.5 + col * 6.45
        y = 1.5 + row * 1.95
        box(sl, x, y, 6.2, 1.8, fill=BG_CARD, border_color=BORDER)
        box(sl, x, y, 0.08, 1.8, fill=color)
        txt(sl, title, x + 0.2, y + 0.12, 5.8, 0.4, size=12, bold=True, color=color)
        txt(sl, body, x + 0.2, y + 0.55, 5.8, 1.1, size=10, color=TEXT_DIM)
        col += 1
        if col == 2:
            col = 0
            row += 1


def slide_market(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4, color=GREEN)
    txt(sl, "Market Opportunity", 0.72, 0.38, 8, 0.6, size=28, bold=True, color=TEXT_WHITE)
    txt(sl, "10", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    # Big TAM number
    box(sl, 0.5, 1.15, 12.33, 1.6, fill=RGBColor(0x06, 0x2a, 0x14), border_color=GREEN)
    txt(sl, "$15B+", 0.7, 1.2, 12.0, 1.0, size=64, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(sl, "annual incident resolution labour + downtime cost (primary TAM)",
        0.7, 2.1, 12.0, 0.45, size=13, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    segments = [
        ("Primary",
         "Enterprise SRE Automation",
         "33,000+ companies with dedicated SRE teams\n"
         "12–40 P1 incidents/month · 40 min saved each\n"
         "$5,600/min Gartner benchmark",
         GREEN, "$15B+"),
        ("Secondary",
         "MSP Platform",
         "One engineer manages 10× more environments\n"
         "White-label per-site licence\n"
         "Agent drafts · engineer approves",
         PURPLE_LT, "$3B+"),
        ("Tertiary",
         "Compliance Verticals",
         "Same framework · different playbooks\n"
         "HIPAA · PCI-DSS · SOX · zero core changes\n"
         "New domain = new registry entries only",
         AMBER, "$2B+"),
    ]

    seg_x = 0.5
    for tier, name, body, col, tam in segments:
        box(sl, seg_x, 2.85, 4.1, 3.8, fill=BG_CARD, border_color=col)
        box(sl, seg_x, 2.85, 4.1, 0.06, fill=col)
        pill(sl, tier, seg_x + 0.15, 2.97, col, BG_MID, size=9, w=1.0, h=0.26)
        txt(sl, name, seg_x + 1.3, 2.95, 2.65, 0.38, size=11, bold=True, color=col)
        txt(sl, body, seg_x + 0.15, 3.35, 3.8, 1.5, size=10, color=TEXT_DIM)
        txt(sl, tam, seg_x + 0.15, 5.55, 3.8, 0.6, size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
        seg_x += 4.25

    txt(sl, "Conservative penetration: 2% of primary TAM at $2K/month = $60M ARR",
        0.5, 6.9, 12.33, 0.45, size=12, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)


def slide_traction(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4, color=GREEN)
    txt(sl, "Traction & Proof Points", 0.72, 0.38, 8, 0.6, size=28, bold=True, color=TEXT_WHITE)
    txt(sl, "11", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    items = [
        ("✅", "AgentCore", "121 tests, 0 failures, no live DB or LLM required to run the full test suite"),
        ("✅", "AgentGateway", "JWT auth + session binding + context budget management shipped and tested"),
        ("✅", "AuthService", "RS256 JWT issuance with JWKS endpoint — production-grade key rotation"),
        ("✅", "Plan Memory", "Full feature implemented and demo'd: backend matching → SSE event → UI badge"),
        ("✅", "React UI", "Live SSE streaming, arc budget gauge, compaction banners, plan history panel"),
        ("✅", "4 Demo Scenarios", "K8s / DB / VPN / SSL — all run in docker compose, no credentials needed"),
        ("✅", "RegistryService", "Full system design + API spec + ScalableRegistryClient (60s TTL) in AgentCore"),
        ("✅", "Previous Hackathon", "HITL + DAG parallel execution demoed and validated with judges"),
        ("✅", "Recorded Demo", "Full MP4 walkthrough: login → clarification → few-shot → HITL → execution → report"),
    ]

    item_y = 1.22
    for icon, title, body in items:
        box(sl, 0.5, item_y, 12.33, 0.65, fill=BG_CARD, border_color=BORDER)
        txt(sl, icon, 0.65, item_y + 0.12, 0.45, 0.38, size=14, color=GREEN)
        txt(sl, title, 1.2, item_y + 0.12, 2.3, 0.38, size=11, bold=True, color=GREEN)
        txt(sl, body, 3.65, item_y + 0.12, 8.95, 0.42, size=10.5, color=TEXT_DIM)
        item_y += 0.73


def slide_ask(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4, color=PURPLE)
    txt(sl, "What We're Raising", 0.72, 0.38, 8, 0.6, size=28, bold=True, color=TEXT_WHITE)
    txt(sl, "12", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    # Amount
    box(sl, 0.5, 1.15, 5.5, 2.2, fill=RGBColor(0x1e, 0x10, 0x45), border_color=PURPLE)
    txt(sl, "$1.5M", 0.7, 1.2, 5.1, 1.2, size=60, bold=True, color=PURPLE_LT, align=PP_ALIGN.CENTER)
    txt(sl, "Pre-Seed · 18-month runway", 0.7, 2.25, 5.1, 0.7, size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Milestone
    box(sl, 6.2, 1.15, 6.6, 2.2, fill=BG_CARD, border_color=BORDER)
    txt(sl, "Milestones", 6.4, 1.25, 6.2, 0.4, size=12, bold=True, color=PURPLE_LT)
    milestones = [
        "20 paid design partners → Series A trigger",
        "$5M ARR target at Series A",
        "RegistryService GA in Year 2",
        "MSP channel launch in Year 2",
    ]
    m_y = 1.72
    for m in milestones:
        txt(sl, f"→  {m}", 6.4, m_y, 6.2, 0.35, size=10.5, color=TEXT_DIM)
        m_y += 0.4

    # Use of funds
    funds = [
        ("55%", "Engineering", "2 hires: backend + frontend", PURPLE_LT),
        ("25%", "Sales & Design Partners", "Outreach, onboarding, design partner program", GREEN),
        ("10%", "Infrastructure", "Cloud, CI/CD, security audits", RGBColor(0x38, 0xbd, 0xf8)),
        ("10%", "Legal / Compliance", "SOC2, legal entity, compliance certs", AMBER),
    ]
    fund_y = 3.6
    txt(sl, "Use of Funds", 0.5, 3.5, 12.33, 0.4, size=13, bold=True, color=TEXT_DIM)
    for pct, title, body, col in funds:
        box(sl, 0.5, fund_y, 12.33, 0.78, fill=BG_CARD, border_color=BORDER)
        box(sl, 0.5, fund_y, 0.08, 0.78, fill=col)
        txt(sl, pct, 0.72, fund_y + 0.18, 0.8, 0.4, size=18, bold=True, color=col)
        txt(sl, title, 1.6, fund_y + 0.08, 3.5, 0.35, size=11, bold=True, color=col)
        txt(sl, body, 1.6, fund_y + 0.44, 10.5, 0.3, size=10, color=TEXT_DIM)
        fund_y += 0.86

    # GTM phases
    box(sl, 0.5, 7.05, 12.33, 0.32, fill=BG_MID, border_color=BORDER)
    txt(sl, "Phase 1: $480K ARR (20 design partners)  ·  Phase 2: $5M ARR (MSP channel)  ·  Phase 3: $25M ARR (RegistryService GA)",
        0.7, 7.1, 11.9, 0.22, size=9, color=TEXT_DIM, align=PP_ALIGN.CENTER)


def slide_why_now(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    accent_bar(sl, 0.5, 0.4, color=AMBER)
    txt(sl, "Why Now", 0.72, 0.38, 8, 0.6, size=28, bold=True, color=TEXT_WHITE)
    txt(sl, "13", 12.0, 0.38, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    tailwinds = [
        ("Claude's Context Window Is the New Production Constraint",
         "As agents handle longer sessions, context budget management is no longer optional — "
         "it's the difference between a working agent and a degraded one. "
         "We built the gateway that manages it. Nobody else has.",
         PURPLE_LT),
        ("Enterprise HITL Is Becoming Regulatory",
         "EU AI Act, NIST AI RMF, and SOC2 Type II all point toward mandatory human oversight "
         "for consequential AI actions. AgenticFrmk's hard HITL gate is a compliance feature, not a UX choice. "
         "Customers are asking for this now.",
         AMBER),
        ("LangGraph Has Crossed the Production Threshold",
         "The primitives — interrupt(), Send API, PostgresSaver — now exist to build production agents. "
         "12 months ago this required custom plumbing. Today the framework is stable enough for enterprise. "
         "The window to establish an infrastructure standard is open now.",
         GREEN),
    ]

    tail_y = 1.3
    for i, (title, body, col) in enumerate(tailwinds):
        box(sl, 0.5, tail_y, 12.33, 1.8, fill=BG_CARD, border_color=col)
        box(sl, 0.5, tail_y, 0.12, 1.8, fill=col)
        txt(sl, str(i + 1), 0.65, tail_y + 0.55, 0.45, 0.6, size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(sl, title, 1.25, tail_y + 0.12, 11.0, 0.45, size=14, bold=True, color=col)
        txt(sl, body, 1.25, tail_y + 0.65, 11.0, 1.0, size=11, color=TEXT_DIM)
        tail_y += 1.95

    box(sl, 0.5, 7.05, 12.33, 0.3, fill=BG_MID)
    txt(sl, "The companies that own the agentic infrastructure layer will own the enterprise AI stack. That window is now.",
        0.7, 7.08, 11.9, 0.26, size=10.5, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)


def slide_team_contact(prs):
    sl = add_slide(prs)
    fill_slide(sl, BG_DARK)

    # Decorative side panels
    box(sl, 0, 0, 0.5, 7.5, fill=RGBColor(0x12, 0x06, 0x2a))
    box(sl, 12.83, 0, 0.5, 7.5, fill=RGBColor(0x12, 0x06, 0x2a))

    txt(sl, "Team & Contact", 1.5, 0.5, 10, 0.65,
        size=30, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "14", 12.0, 0.5, 1.0, 0.5, size=11, color=TEXT_DIMMER, align=PP_ALIGN.RIGHT)

    box(sl, 0.8, 1.4, 11.73, 3.0, fill=BG_CARD, border_color=BORDER)
    txt(sl, "Founder", 1.0, 1.55, 11.0, 0.45, size=11, color=TEXT_DIMMER, bold=True)
    txt(sl, "Full-stack systems engineer", 1.0, 2.05, 11.0, 0.55, size=20, bold=True, color=TEXT_WHITE)
    txt(sl,
        "Background in distributed systems, LLM infrastructure, and enterprise SRE tooling.\n"
        "Built AgenticFrmk end-to-end: AgentCore (121 tests), AgentGateway, AuthService,\n"
        "RegistryService design, SREDemo React UI, Docker stack, and complete test coverage.\n"
        "Won previous Agentic Seattle hackathon — HITL + DAG parallel execution.",
        1.0, 2.65, 11.3, 1.5, size=12, color=TEXT_DIM)

    divider(sl, 4.65)

    contact_items = [
        ("Email",  "mauttaram@gmail.com",           PURPLE_LT),
        ("GitHub", "github.com/AgenticFrmk",        RGBColor(0x38, 0xbd, 0xf8)),
        ("Demo",   "docker compose up → localhost:3000", GREEN),
    ]
    c_x = 1.5
    for label, value, col in contact_items:
        box(sl, c_x, 4.85, 3.5, 1.4, fill=BG_CARD, border_color=col)
        txt(sl, label, c_x + 0.15, 4.95, 3.2, 0.38, size=10, color=TEXT_DIMMER, bold=True)
        txt(sl, value, c_x + 0.15, 5.38, 3.2, 0.65, size=11, color=col)
        c_x += 3.6

    txt(sl, "AgenticFrmk", 1.5, 6.45, 10, 0.6,
        size=30, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "Production-grade agentic infrastructure · Built to deploy",
        1.5, 7.05, 10, 0.38, size=13, color=PURPLE_LT, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def build() -> Path:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_problem(prs)
    slide_solution_overview(prs)
    slide_gateway(prs)
    slide_plan_memory(prs)
    slide_hitl_registry(prs)
    slide_sre_demo_value(prs)
    slide_architecture(prs)
    slide_moat(prs)
    slide_market(prs)
    slide_traction(prs)
    slide_ask(prs)
    slide_why_now(prs)
    slide_team_contact(prs)

    out = Path(__file__).parent / "AgenticFrmk-VC-Pitch.pptx"
    prs.save(str(out))
    print(f"✅  Saved: {out}  ({len(prs.slides)} slides)")
    return out


if __name__ == "__main__":
    build()
